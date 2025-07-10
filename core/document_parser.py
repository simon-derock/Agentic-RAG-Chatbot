"""Document Parser for Multi-Format Support"""

import io
import PyPDF2
from pptx import Presentation
import pandas as pd
from docx import Document
import markdown
from typing import List, Dict, Any, Tuple

class DocumentParser:
    """Unified document parsing for multiple formats"""
    
    @staticmethod
    def parse_pdf(file_data: bytes) -> List[Dict[str, Any]]:
        """Parse PDF and extract text chunks"""
        chunks = []
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_data))
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    chunks.append({
                        "text": text.strip(),
                        "page": page_num,
                        "type": "pdf"
                    })
        except Exception as e:
            raise ValueError(f"PDF parsing error: {str(e)}")
        return chunks
    
    @staticmethod
    def parse_pptx(file_data: bytes) -> List[Dict[str, Any]]:
        """Parse PPTX and extract slide content"""
        chunks = []
        try:
            prs = Presentation(io.BytesIO(file_data))
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    chunks.append({
                        "text": "\n".join(slide_text),
                        "slide": slide_num,
                        "type": "pptx"
                    })
        except Exception as e:
            raise ValueError(f"PPTX parsing error: {str(e)}")
        return chunks
    
    @staticmethod
    def parse_csv(file_data: bytes) -> List[Dict[str, Any]]:
        """Parse CSV and convert to text chunks"""
        chunks = []
        try:
            df = pd.read_csv(io.StringIO(file_data.decode('utf-8')))
            
            # Header chunk
            chunks.append({
                "text": f"CSV Headers: {', '.join(df.columns.tolist())}",
                "row": 0,
                "type": "csv"
            })
            
            # Row chunks (group by 10 rows for efficiency)
            for i in range(0, len(df), 10):
                chunk_df = df.iloc[i:i+10]
                text_repr = chunk_df.to_string(index=False)
                chunks.append({
                    "text": text_repr,
                    "row": i+1,
                    "type": "csv"
                })
        except Exception as e:
            raise ValueError(f"CSV parsing error: {str(e)}")
        return chunks
    
    @staticmethod
    def parse_docx(file_data: bytes) -> List[Dict[str, Any]]:
        """Parse DOCX and extract paragraphs"""
        chunks = []
        try:
            doc = Document(io.BytesIO(file_data))
            for para_num, paragraph in enumerate(doc.paragraphs, 1):
                if paragraph.text.strip():
                    chunks.append({
                        "text": paragraph.text.strip(),
                        "paragraph": para_num,
                        "type": "docx"
                    })
        except Exception as e:
            raise ValueError(f"DOCX parsing error: {str(e)}")
        return chunks
    
    @staticmethod
    def parse_text(file_data: bytes, file_extension: str) -> List[Dict[str, Any]]:
        """Parse TXT/MD files"""
        chunks = []
        try:
            text = file_data.decode('utf-8')
            
            if file_extension == '.md':
                # Convert markdown to HTML then extract text
                html = markdown.markdown(text)
                text = html
            
            # Split by paragraphs
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            for para_num, paragraph in enumerate(paragraphs, 1):
                chunks.append({
                    "text": paragraph,
                    "paragraph": para_num,
                    "type": file_extension[1:]  # Remove dot
                })
        except Exception as e:
            raise ValueError(f"Text parsing error: {str(e)}")
        return chunks
    
    @classmethod
    def parse_document(cls, file_data: bytes, file_name: str) -> List[Dict[str, Any]]:
        """Main parsing method - routes to appropriate parser"""
        file_extension = file_name.lower().split('.')[-1]
        
        parser_map = {
            'pdf': cls.parse_pdf,
            'pptx': cls.parse_pptx,
            'csv': cls.parse_csv,
            'docx': cls.parse_docx,
            'txt': lambda data: cls.parse_text(data, '.txt'),
            'md': lambda data: cls.parse_text(data, '.md')
        }
        
        if file_extension not in parser_map:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        chunks = parser_map[file_extension](file_data)
        
        # Add source metadata to all chunks
        for chunk in chunks:
            chunk['source'] = file_name
        
        return chunks